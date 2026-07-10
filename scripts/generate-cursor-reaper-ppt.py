#!/usr/bin/env python3
"""Generate Cursor vs Reaper comparison PowerPoint."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "Cursor-vs-Reaper-Comparison.pptx"

# Brand-ish palette
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xE8, 0x4D, 0x3D)  # reaper red
ACCENT2 = RGBColor(0x3D, 0x9E, 0xE8)  # cursor blue
TEXT = RGBColor(0xF5, 0xF5, 0xF5)
MUTED = RGBColor(0xB0, 0xB0, 0xC0)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
AMBER = RGBColor(0xFF, 0xB3, 0x00)


def set_slide_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(18)
    sp.font.color.rgb = MUTED
    sp.alignment = PP_ALIGN.CENTER

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.4))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "Reaper v0.1.3 · July 2026"
    fp.font.size = Pt(12)
    fp.font.color.rgb = MUTED
    fp.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.1), Inches(10), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(8.8), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER


def add_bullet_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    hdr = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.7))
    hp = hdr.text_frame.paragraphs[0]
    hp.text = title
    hp.font.size = Pt(28)
    hp.font.bold = True
    hp.font.color.rgb = TEXT

    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(9), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, color = item
            p.text = text
            p.font.color.rgb = color
        else:
            p.text = item
            p.font.color.rgb = TEXT
        p.font.size = Pt(16)
        p.space_after = Pt(8)
        p.level = 0

    if note:
        nb = slide.shapes.add_textbox(Inches(0.55), Inches(6.5), Inches(9), Inches(0.5))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(11)
        np.font.color.rgb = MUTED


def add_comparison_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    hdr = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    hp = hdr.text_frame.paragraphs[0]
    hp.text = title
    hp.font.size = Pt(26)
    hp.font.bold = True
    hp.font.color.rgb = TEXT

    cols = len(headers)
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.4), Inches(1.05), Inches(9.2), Inches(5.8)).table

    col_widths = [2.2, 3.5, 3.5]
    for ci, w in enumerate(col_widths[:cols]):
        table.columns[ci].width = Inches(w)

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x44)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT if ci == 0 else (ACCENT2 if ci == 1 else ACCENT)
            p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x35) if ri % 2 else RGBColor(0x1E, 0x1E, 0x30)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT
                if ci > 0:
                    p.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    hdr = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    hp = hdr.text_frame.paragraphs[0]
    hp.text = title
    hp.font.size = Pt(26)
    hp.font.bold = True
    hp.font.color.rgb = TEXT

    for side, (stitle, items, x) in enumerate([
        (left_title, left_items, 0.45),
        (right_title, right_items, 5.05),
    ]):
        color = ACCENT2 if side == 0 else ACCENT
        th = slide.shapes.add_textbox(Inches(x), Inches(1.0), Inches(4.4), Inches(0.45))
        tp = th.text_frame.paragraphs[0]
        tp.text = stitle
        tp.font.size = Pt(18)
        tp.font.bold = True
        tp.font.color.rgb = color

        body = slide.shapes.add_textbox(Inches(x), Inches(1.45), Inches(4.4), Inches(5.2))
        tf = body.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT
            p.space_after = Pt(6)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Cursor vs Reaper",
        "Feature comparison for AI-assisted development workflows",
    )

    add_bullet_slide(
        prs,
        "Executive summary",
        [
            ("Cursor", ACCENT2) if False else "Cursor — AI-native IDE (VS Code fork) built around autonomous agents, codebase indexing, and cloud execution.",
            "Reaper — Local developer git studio (Rust) with IntelliJ-style UI, self-hosted repos, and multi-provider AI (Cursor bridge + Gemini).",
            "Cursor optimizes for: broad language support, agent autonomy, team/cloud workflows, extension ecosystem.",
            "Reaper optimizes for: git-centric local dev, Java/enterprise builds, private PAT remotes, integrated build/DB/coverage panels.",
            "They overlap on: inline AI completions, agent chat with file edits, terminal, git UI — but serve different primary use cases.",
        ],
        note="Audience: engineering leads evaluating tooling for Java/enterprise teams vs general AI-first development.",
    )

    add_two_column_slide(
        prs,
        "Product positioning",
        "Cursor",
        [
            "VS Code fork → full extension marketplace",
            "Agent-first interface (Cursor 3)",
            "Semantic codebase index + @-mentions",
            "Cloud & background agents (PR output)",
            "Credit-metered SaaS ($20–200/mo)",
            "macOS, Windows, Linux + mobile + CLI",
        ],
        "Reaper",
        [
            "Rust native app + Monaco editor",
            "Git studio: host bare repos over HTTP",
            "IntelliJ-inspired single-window UI",
            "Local-first; PAT-based private remotes",
            "Open source (MIT); self-hosted data",
            "macOS DMG (server mode cross-platform)",
        ],
    )

    add_section_slide(prs, "AI & agents")

    add_comparison_table_slide(
        prs,
        "AI capabilities",
        ["Capability", "Cursor", "Reaper"],
        [
            ("Chat modes", "Ask · Agent · Manual/Edit", "Cursor: Agent/Plan/Ask · Gemini: read-only Q&A"),
            ("File edits", "Agent + Composer multi-file diffs", "Cursor bridge edits + revert last turn"),
            ("Inline completions", "Tab completions (model picker)", "Cursor → Gemini → Claude → LSP chain"),
            ("Cloud agents", "Background VMs, async PRs", "Not available (local only)"),
            ("Parallel agents", "Up to 8 (worktrees / cloud)", "Single agent panel"),
            ("Context", "@codebase @file @docs @web", "Active file + repo; no @-mention UX"),
            ("Rules / skills", ".cursor/rules, skills, MCP plugins", "Not available"),
            ("PR review", "BugBot automated review", "Not available"),
            ("Commit messages", "Via agent/chat", "Gemini suggest from staged diff"),
            ("Quick fixes", "Via agent / inline", "Gemini + Cursor fallback bulb"),
        ],
    )

    add_section_slide(prs, "Editor & IDE")

    add_comparison_table_slide(
        prs,
        "Editor & language support",
        ["Area", "Cursor", "Reaper"],
        [
            ("Editor core", "VS Code + extensions", "Monaco + custom providers"),
            ("Themes / UX", "VS Code themes", "6 themes, IntelliJ-style layout"),
            ("LSP depth", "Full VS Code LSP ecosystem", "Deep: Java (jdtls), C/C++ (clangd)"),
            ("Other languages", "TS, Python, Rust, Go, etc.", "Syntax + formatters; limited LSP"),
            ("Java tooling", "Via extensions", "Classpath index, javac diag, Spring, JaCoCo"),
            ("Debugger", "Full VS Code debugger", "Not available"),
            ("Refactoring", "LSP + AI agent", "Rename, Find Usages, Format, Organize Imports"),
            ("Inline ghost text", "Yes (Copilot-style)", "Yes (Tabnine-style, 1294 regression tests)"),
            ("Split editor", "Yes", "Not available"),
            ("Regression tests", "External / user", "1294+ tests on every cargo build"),
        ],
    )

    add_section_slide(prs, "Git & DevOps")

    add_comparison_table_slide(
        prs,
        "Git & development workflow",
        ["Area", "Cursor", "Reaper"],
        [
            ("Git UI", "Source control panel (VS Code)", "Commit panel, log, branch picker, conflicts"),
            ("Self-hosted git", "Uses your remotes only", "Hosts bare repos + smart HTTP"),
            ("Private enterprise", "Clone any remote", "PAT import: GitHub, GitLab, Azure, etc."),
            ("Publish workflow", "Via git remote / GitHub", "Built-in Publish to GitHub"),
            ("Secret scanning", "Not built-in", "Pre-commit / pre-push scan"),
            ("Merge conflicts", "VS Code merge editor", "Per-hunk Ours/Theirs/Both UI"),
            ("Git terminal", "Integrated terminal", "Whitelisted git console + xterm"),
            ("Build tasks", "Tasks / terminal", "Maven/Gradle/CMake/npm tree panel"),
            ("Coverage", "Via extensions", "JaCoCo inline + panel (Java)"),
            ("Database viewer", "Via MCP / extensions", "Built-in Postgres + SQLite"),
            ("Docker logs", "Terminal / extensions", "Built-in compose log stream"),
        ],
    )

    add_section_slide(prs, "Strengths & gaps")

    add_two_column_slide(
        prs,
        "Where each product wins",
        "Cursor strengths",
        [
            "Best-in-class agent autonomy & cloud execution",
            "Entire VS Code extension + LSP ecosystem",
            "Codebase-wide semantic index & @-context",
            "MCP plugins, rules, skills, BugBot",
            "Multi-platform + team billing & admin",
            "Ideal for polyglot teams & greenfield AI workflows",
        ],
        "Reaper strengths",
        [
            "Self-hosted git studio (bare repo HTTP hosting)",
            "Enterprise PAT import without leaving the app",
            "Java-first: Spring Boot, JaCoCo, jdtls, javac loop",
            "Cursor API bridge — agent models without Cursor IDE",
            "Integrated build tasks, DB viewer, Docker logs",
            "Local data ownership; MIT open source",
        ],
    )

    add_bullet_slide(
        prs,
        "Reaper gaps vs Cursor",
        [
            "No debugger (breakpoints, watch, call stack)",
            "No extension marketplace or MCP plugin system",
            "Claude agent slot is placeholder (inline only)",
            "Limited LSP beyond Java and C/C++",
            "No cloud/background agents or BugBot",
            "No @codebase context UI or project rules",
            "macOS GUI primary; no Windows/Linux installer yet",
            "No split editor, keybinding editor, or SSH remote",
        ],
        note="Many gaps are intentional — Reaper targets git-centric local Java/enterprise workflows, not general IDE replacement.",
    )

    add_bullet_slide(
        prs,
        "Cursor gaps vs Reaper",
        [
            "No self-hosted bare git repository hosting",
            "No built-in secret scanning on commit/push",
            "No integrated JaCoCo coverage or Java classpath index",
            "No built-in database viewer or Docker log panel",
            "Requires Cursor subscription for full agent features",
            "Less depth for Maven/Gradle/Spring Boot run workflows",
            "No agent revert per-turn for file edits",
            "Vendor lock-in to Anysphere cloud for background agents",
        ],
    )

    add_section_slide(prs, "Decision guide")

    add_comparison_table_slide(
        prs,
        "When to choose which",
        ["Scenario", "Better fit", "Why"],
        [
            ("Polyglot startup, AI-first culture", "Cursor", "Agents, extensions, cloud PRs"),
            ("Java/Spring enterprise monorepo", "Reaper", "jdtls, JaCoCo, build tasks, Spring run"),
            ("Self-hosted / air-gapped git studio", "Reaper", "Local ~/reaper, bare HTTP git"),
            ("Private GitLab/Azure with PATs", "Reaper", "Built-in PAT import & sync"),
            ("Multi-repo parallel cloud agents", "Cursor", "Cursor 3 agent workspace"),
            ("Use Cursor models, own your IDE", "Reaper", "Cursor bridge without Cursor app"),
            ("Full debugger + test explorer", "Cursor", "VS Code debugging ecosystem"),
            ("Commit hygiene + secret prevention", "Reaper", "Secret scan + AI commit msg"),
        ],
    )

    add_bullet_slide(
        prs,
        "Hybrid strategy",
        [
            "Teams can use Reaper for day-to-day Java/git work and Cursor for agent-heavy refactors.",
            "Reaper's Cursor bridge lets developers consume Cursor agent API without switching IDEs.",
            "Gemini in Reaper covers read-only Q&A, commit messages, and inline completions at lower cost.",
            "Reaper's regression suite (1294 tests) guards editor behavior on every build — unique for a git studio.",
            "Roadmap opportunity: Claude agent, rules/skills, debugger — closing gaps while keeping git-studio focus.",
        ],
    )

    add_bullet_slide(
        prs,
        "Summary",
        [
            ("Cursor = AI-native general-purpose IDE with cloud agents", ACCENT2),
            ("Reaper = local git studio with Java depth + multi-provider AI", ACCENT),
            "Not direct substitutes — complementary for enterprise Java teams.",
            "Reaper differentiates on self-hosted git, PAT workflows, build/DB/coverage integration.",
            "Cursor differentiates on agent autonomy, ecosystem breadth, and team/cloud scale.",
            "Evaluate based on: language mix, hosting model, AI budget, and compliance requirements.",
        ],
    )

    add_title_slide(
        prs,
        "Questions?",
        "docs/Cursor-vs-Reaper-Comparison.pptx · reaper-org/reaper",
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
